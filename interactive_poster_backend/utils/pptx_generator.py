import pptx
from typing import Optional, Union
from typing import Optional, Union
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE # Not used in current logic, but can be kept for future
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import logging
import requests
from io import BytesIO
from pathlib import Path # Ensure Path is imported
from .. import config # Ensure config is imported

from ..schemas import models as schemas

logger = logging.getLogger(__name__)

def apply_theme_to_slide(slide, selected_theme: str, style_overrides: Optional[schemas.PosterElementStyles]):
    """Applies background color based on the theme, then applies overrides."""

    # Default to white or transparent depending on slide master
    # Forcing solid white for themes that don't specify a dark background
    # to ensure text visibility if text becomes light due to theme/override.
    slide_fill = slide.background.fill
    slide_fill.solid()
    if selected_theme in ["default", "professional_blue"]:
        slide_fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White

    # Theme-based background
    if selected_theme == "minimalist_dark":
        background = slide.background # Already got fill above
        slide_fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E) # Very dark grey
    elif selected_theme == "creative_warm":
        # background = slide.background # fill is already slide_fill
        slide_fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xF0) # Light creamy background

    # Apply slide_background override if present
    if style_overrides and style_overrides.slide_background:
        try:
            hex_color = style_overrides.slide_background.lstrip('#')
            if len(hex_color) == 6: # Basic validation
                slide_fill.solid() # Ensure it's solid fill before setting color
                slide_fill.fore_color.rgb = RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
            else:
                logger.warning(f"Invalid hex color format for slide_background override: '{style_overrides.slide_background}'")
        except ValueError as e:
            logger.warning(f"Invalid slide_background color value '{style_overrides.slide_background}': {e}")
        except Exception as e: # Catch any other parsing errors
            logger.error(f"Error applying slide_background override '{style_overrides.slide_background}': {e}")


def apply_font_style(text_frame, role: str, selected_theme: str,
                     style_overrides: Optional[schemas.PosterElementStyles],
                     is_body: bool = False):
    """Applies font color, size, and family based on theme and then overrides."""
    if not text_frame.paragraphs:
        return

    # Ensure the first paragraph has at least one run to attach font styles to
    # This is crucial if the text_frame was empty and text was just set.
    paragraph = text_frame.paragraphs[0]
    if not paragraph.runs:
        run = paragraph.add_run()
        run.text = paragraph.text # If text was on paragraph, move to run
        paragraph.text = "" # Clear text from paragraph itself if it was moved

    font = paragraph.font # Get font of the first paragraph for theme styling
    # If there are specific runs and we want to style them all, we'd iterate paragraph.runs

    # --- Step 1: Apply Theme-Based Styles ---
    # Default text color for "default" theme (usually black, but can be set explicitly)
    if selected_theme == "default":
        font.color.rgb = RGBColor(0x10, 0x10, 0x10) # Dark grey/near black

    if selected_theme == "minimalist_dark":
        font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        if is_body:
            font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    elif selected_theme == "professional_blue":
        if role == "title": font.color.rgb = RGBColor(0x00, 0x33, 0x66)
        elif role == "section_title": font.color.rgb = RGBColor(0x00, 0x5A, 0x9C)
        else: font.color.rgb = RGBColor(0x22, 0x22, 0x22) # Body
    elif selected_theme == "creative_warm":
        if role == "title": font.color.rgb = RGBColor(0xB7, 0x41, 0x0E)
        elif role == "section_title": font.color.rgb = RGBColor(0xD9, 0x53, 0x4F)
        else: font.color.rgb = RGBColor(0x3C, 0x3C, 0x3C) # Body

    # Theme-based font size adjustments
    if role == "title": font.size = Pt(40); paragraph.alignment = PP_ALIGN.CENTER
    elif role == "subtitle": font.size = Pt(24); paragraph.alignment = PP_ALIGN.CENTER
    elif role == "section_title": font.size = Pt(32)
    elif is_body: font.size = Pt(16)

    # --- Step 2: Apply Overrides ---
    override_props: Optional[schemas.ElementStyleProperties] = None
    if style_overrides:
        if role == "title" and style_overrides.title: override_props = style_overrides.title
        elif role == "subtitle" and style_overrides.abstract: override_props = style_overrides.abstract # Assuming subtitle maps to abstract style
        elif role == "conclusion_title" and style_overrides.conclusion: override_props = style_overrides.conclusion # For title of conclusion section
        elif role == "conclusion_body" and style_overrides.conclusion: override_props = style_overrides.conclusion # For body of conclusion section
        elif role == "section_title" and style_overrides.section_title: override_props = style_overrides.section_title
        elif is_body and style_overrides.section_content: override_props = style_overrides.section_content # General for all section bodies

    if override_props:
        if override_props.color:
            try:
                hex_color = override_props.color.lstrip('#')
                if len(hex_color) == 6:
                    font.color.rgb = RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
                else:
                    logger.warning(f"Invalid hex color format for role '{role}': '{override_props.color}'")
            except ValueError as e:
                 logger.warning(f"Invalid color value for role '{role}' ('{override_props.color}'): {e}")
            except Exception as e:
                 logger.error(f"Error applying color override for role '{role}' ('{override_props.color}'): {e}")

        if override_props.font_size:
            font.size = Pt(override_props.font_size)
        if override_props.font_family: # Use with caution - font must be available on system viewing PPTX
            try:
                font.name = override_props.font_family
            except Exception as e:
                 logger.warning(f"Failed to set font family '{override_props.font_family}' for role '{role}': {e}")


def generate_pptx_from_data(poster_data: schemas.Poster, output_path: str) -> None:
    prs = pptx.Presentation()
    selected_theme = poster_data.selected_theme if poster_data.selected_theme else "default"
    style_overrides = poster_data.style_overrides # This is Optional[PosterElementStyles]

    # Slide Layouts
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1] # Typically: Title and Content
    # blank_slide_layout = prs.slide_layouts[5] # Blank, if needed for more custom layouts

    # --- Title Slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    apply_theme_to_slide(slide, selected_theme)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1] # Assuming this is the subtitle placeholder

    title_shape.text = poster_data.title
    apply_font_style(title_shape.text_frame, "title", selected_theme)

    if poster_data.abstract:
        subtitle_shape.text = poster_data.abstract
        apply_font_style(subtitle_shape.text_frame, "subtitle", selected_theme)
    else: # Remove subtitle shape if no abstract
        if subtitle_shape.has_text_frame: # Check if it's a placeholder that can be deleted
             sp = subtitle_shape.element
             sp.getparent().remove(sp)


    # --- Sections Slides ---
    for section in poster_data.sections:
        slide = prs.slides.add_slide(content_slide_layout)
        apply_theme_to_slide(slide, selected_theme)

        title_placeholder = slide.shapes.title # Or slide.placeholders[0] if more reliable
        body_placeholder = slide.placeholders[1] # Assuming content placeholder index

        title_placeholder.text = section.section_title
        apply_font_style(title_placeholder.text_frame, "section_title", selected_theme)

        tf_body = body_placeholder.text_frame
        tf_body.clear()
        p_body = tf_body.add_paragraph()
        p_body.text = section.section_content or ""
        apply_font_style(tf_body, "body", selected_theme, style_overrides, is_body=True) # Role "body" for section content

        # Image Embedding Logic
        if section.image_urls: # Now uses 'image_urls'
            logger.info(f"Section '{section.section_title}' has {len(section.image_urls)} image URLs.")

            # These are example coordinates and sizes, adjust based on typical slide master.
            img_left = Inches(5.5) # Example: Right half of a 10-inch wide slide
            img_top = Inches(1.8)  # Example: Below title, allowing space for text on left
            img_max_width = Inches(4.0)
            img_max_height = Inches(3.5)

            image_added_count = 0
            for img_path_or_url in section.image_urls:
                if image_added_count >= 1: # Still only adding the first successfully processed image
                    logger.info(f"Processed first image for section '{section.section_title}'. Skipping further images.")
                    break

                image_source_for_pptx: Union[BytesIO, Path, str]
                is_local_file_from_upload = False # Flag to manage BytesIO closing

                if img_path_or_url.startswith(('http://', 'https://')):
                    try:
                        logger.info(f"Downloading image from URL: {img_path_or_url}")
                        response = requests.get(img_path_or_url, stream=True, timeout=10)
                        response.raise_for_status()
                        image_source_for_pptx = BytesIO(response.content)
                        logger.info(f"Successfully downloaded image from {img_path_or_url}.")
                    except requests.exceptions.RequestException as e_req:
                        logger.error(f"Failed to download image {img_path_or_url} for section '{section.section_title}': {e_req}")
                        continue
                    except Exception as e_general_web:
                        logger.error(f"An unexpected error occurred with web image {img_path_or_url} for section '{section.section_title}': {e_general_web}")
                        continue
                elif img_path_or_url: # Assuming it's a relative local path
                    # Relative path is like "poster_{pid}/section_{sid}/filename.ext"
                    # It's relative to config.UPLOADED_IMAGES_DIR_NAME, which is a subdirectory of APP_ROOT_DIR
                    # config.UPLOADED_IMAGES_DIR gives the absolute path to "uploaded_images"
                    absolute_image_path = config.UPLOADED_IMAGES_DIR / img_path_or_url

                    if not absolute_image_path.is_file():
                        logger.error(f"Local image file not found at resolved path: {absolute_image_path} (original relative: {img_path_or_url})")
                        continue
                    image_source_for_pptx = str(absolute_image_path) # add_picture can take a path string
                    is_local_file_from_upload = True
                    logger.info(f"Found local image file: {absolute_image_path}")
                else: # Empty or None URL
                    logger.warning(f"Empty or invalid image path/URL for section '{section.section_title}'.")
                    continue

                # Common logic for adding picture to slide
                try:
                    pic = slide.shapes.add_picture(image_source_for_pptx, img_left, img_top, width=img_max_width)

                    if pic.height > img_max_height:
                        scale_ratio = img_max_height / pic.height
                        pic.width = int(pic.width * scale_ratio)
                        pic.height = img_max_height

                    # Optional: Further check for slide width overflow if necessary
                    # if pic.left + pic.width > prs.slide_width: ...

                    logger.info(f"Successfully added image '{img_path_or_url}' to section '{section.section_title}'.")
                    image_added_count += 1
                except FileNotFoundError as e_fnf:
                     logger.error(f"Local image file not found by python-pptx at: {image_source_for_pptx}. Error: {e_fnf}")
                except Exception as e_add_pic:
                    logger.error(f"Failed to add image '{img_path_or_url}' to slide for section '{section.section_title}': {e_add_pic}", exc_info=True)
                finally:
                    if not is_local_file_from_upload and isinstance(image_source_for_pptx, BytesIO):
                        image_source_for_pptx.close() # Close BytesIO stream for downloaded images

        elif section.image_urls is None: # Explicitly None
            pass # No images to process
        # else: [] (empty list) is handled by 'if section.image_urls:' failing

    # --- Conclusion Slide ---
    if poster_data.conclusion:
        slide = prs.slides.add_slide(content_slide_layout)
        apply_theme_to_slide(slide, selected_theme)

        title_placeholder = slide.shapes.title
        body_placeholder = slide.placeholders[1]

        title_placeholder.text = "Conclusion"
        apply_font_style(title_placeholder.text_frame, "section_title", selected_theme) # Use section_title style for "Conclusion"

        tf_body = body_placeholder.text_frame
        tf_body.clear()
        p_body = tf_body.add_paragraph()
        p_body.text = poster_data.conclusion
        apply_font_style(tf_body, "body", selected_theme, is_body=True)

    prs.save(output_path)
